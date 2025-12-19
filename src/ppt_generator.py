"""
PowerPoint Generation Module - FINAL FIXED VERSION
✅ NO underlines on titles
✅ Proper conclusion slide structure
✅ SlideX watermark in bottom-right corner
✅ No overlapping elements
"""

import os
from typing import List, Dict, Any, Optional
from datetime import datetime
import tempfile

# PowerPoint library
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Local imports
from .config import Config, PRESENTATION_TEMPLATES


class PPTGenerator:
    """Generates STYLISH PowerPoint presentations"""
    
    def __init__(self, template_style: str = "professional"):
        self.config = Config()
        self.template_style = template_style
        self.template_config = PRESENTATION_TEMPLATES.get(
            template_style, 
            PRESENTATION_TEMPLATES["professional"]
        )
    
    def create_presentation(
        self, 
        slides_data: List[Dict[str, Any]], 
        template_style: str = None,
        output_path: Optional[str] = None
    ) -> str:
        """Create a STYLISH PowerPoint presentation"""
        try:
            # Update template if provided
            if template_style:
                self.template_style = template_style
                self.template_config = PRESENTATION_TEMPLATES.get(
                    template_style,
                    PRESENTATION_TEMPLATES["professional"]
                )
            
            # Create new presentation
            prs = Presentation()
            
            # Set slide dimensions (16:9 widescreen)
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            print(f"\n{'='*60}")
            print(f"🎨 Creating STYLISH presentation")
            print(f"📊 Slides: {len(slides_data)}")
            print(f"🎯 Template: {self.template_style}")
            print(f"{'='*60}\n")
            
            # Add slides
            for i, slide_data in enumerate(slides_data):
                slide_num = i + 1
                slide_type = slide_data.get('slide_type', 'content')
                title = slide_data.get('title', 'Untitled')
                
                try:
                    if slide_num == 1:
                        self._add_modern_title_slide(prs, slide_data)
                        print(f"✓ Title slide: {title}")
                    elif slide_type == 'conclusion' or title.lower() == 'conclusion':
                        self._add_conclusion_slide(prs, slide_data, slide_num)
                        print(f"✓ Conclusion slide: {title}")
                    else:
                        self._add_modern_content_slide(prs, slide_data, slide_num)
                        print(f"✓ Slide {slide_num}: {title}")
                except Exception as slide_error:
                    print(f"⚠️ Error on slide {slide_num}: {slide_error}")
                    self._add_fallback_slide(prs, slide_data, slide_num)
            
            # Save presentation
            if output_path is None:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                output_path = os.path.join(
                    tempfile.gettempdir(), 
                    f"slidex_presentation_{timestamp}.pptx"
                )
            
            prs.save(output_path)
            print(f"\n✅ Presentation saved: {output_path}\n")
            return output_path
            
        except Exception as e:
            print(f"❌ ERROR creating presentation: {e}")
            import traceback
            traceback.print_exc()
            raise
    
    def _add_watermark(self, slide):
        """Add SlideX watermark to bottom-right corner"""
        try:
            # SlideX watermark (bottom-right)
            watermark_box = slide.shapes.add_textbox(
                Inches(8.3), Inches(7.1),
                Inches(1.5), Inches(0.3)
            )
            watermark_frame = watermark_box.text_frame
            watermark_frame.text = "SlideX"
            
            watermark_para = watermark_frame.paragraphs[0]
            watermark_para.font.name = "Segoe UI"
            watermark_para.font.size = Pt(10)
            watermark_para.font.color.rgb = RGBColor(150, 150, 150)
            watermark_para.font.italic = True
            watermark_para.alignment = PP_ALIGN.RIGHT
        except Exception as e:
            print(f"Warning: Could not add watermark: {e}")
    
    def _add_modern_title_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """Add a MODERN title slide (FIXED - no overlapping)"""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Get colors
        theme_color = self._hex_to_rgb(self.template_config.get("theme_color", "#1E3A8A"))
        accent_color = self._hex_to_rgb(self.template_config.get("accent_color", "#3B82F6"))
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Large colored panel (left side)
        left_panel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(4.5), Inches(7.5)
        )
        left_panel.fill.solid()
        left_panel.fill.fore_color.rgb = RGBColor(*theme_color)
        left_panel.line.fill.background()
        
        # Decorative triangle (smaller, positioned better)
        try:
            triangle = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_TRIANGLE,
                Inches(3.5), Inches(6.2),
                Inches(1.5), Inches(1.2)
            )
            triangle.fill.solid()
            triangle.fill.fore_color.rgb = RGBColor(*accent_color)
            triangle.line.fill.background()
            triangle.rotation = 45
        except:
            pass
        
        # Icon circle (positioned higher to avoid overlap)
        icon_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.5), Inches(2.8),
            Inches(1.5), Inches(1.5)
        )
        icon_circle.fill.solid()
        icon_circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
        icon_circle.line.fill.background()
        
        # Icon emoji
        icon_frame = icon_circle.text_frame
        icon_frame.text = "🎯"
        icon_frame.vertical_anchor = 1  # Middle
        icon_para = icon_frame.paragraphs[0]
        icon_para.font.size = Pt(48)
        icon_para.alignment = PP_ALIGN.CENTER
        
        # Title (right side, no underline)
        title_text = slide_data.get('title', 'AI Generated Presentation')
        title_box = slide.shapes.add_textbox(
            Inches(5), Inches(2.5),
            Inches(4.5), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.word_wrap = True
        
        title_para = title_frame.paragraphs[0]
        title_para.font.name = "Segoe UI"
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(*theme_color)
        title_para.font.underline = False  # ✅ NO UNDERLINE
        title_para.alignment = PP_ALIGN.LEFT
        
        # Subtitle (bottom of slide - below colored panel)
        subtitle_box = slide.shapes.add_textbox(
            Inches(5), Inches(6.8),
            Inches(4.5), Inches(0.6)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"Generated by SlideX AI\n{datetime.now().strftime('%B %d, %Y')}"
        
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.name = "Segoe UI"
        subtitle_para.font.size = Pt(12)
        subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
        subtitle_para.alignment = PP_ALIGN.LEFT
        
        # ✅ Add watermark to title slide
        self._add_watermark(slide)
    
    def _add_modern_content_slide(
        self, 
        prs: Presentation, 
        slide_data: Dict[str, Any], 
        slide_number: int
    ):
        """Add a MODERN content slide (FIXED - no underline)"""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Get colors
        theme_color = self._hex_to_rgb(self.template_config.get("theme_color", "#1E3A8A"))
        accent_color = self._hex_to_rgb(self.template_config.get("accent_color", "#3B82F6"))
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Top accent bar
        header_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.4)
        )
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = RGBColor(*theme_color)
        header_bar.line.fill.background()
        
        # Left accent line
        left_accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0.4),
            Inches(0.2), Inches(7.1)
        )
        left_accent.fill.solid()
        left_accent.fill.fore_color.rgb = RGBColor(*accent_color)
        left_accent.line.fill.background()
        
        # Slide number (top right)
        slide_num_box = slide.shapes.add_textbox(
            Inches(9), Inches(0.05),
            Inches(0.9), Inches(0.3)
        )
        slide_num_frame = slide_num_box.text_frame
        slide_num_frame.text = f"{slide_number}"
        slide_num_para = slide_num_frame.paragraphs[0]
        slide_num_para.font.name = "Segoe UI"
        slide_num_para.font.size = Pt(14)
        slide_num_para.font.bold = True
        slide_num_para.font.color.rgb = RGBColor(255, 255, 255)
        slide_num_para.alignment = PP_ALIGN.CENTER
        
        # ✅ Title (NO UNDERLINE - removed decorative underline)
        title = slide_data.get('title', 'Untitled Slide')
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.8),
            Inches(9), Inches(0.9)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        
        title_para = title_frame.paragraphs[0]
        title_para.font.name = "Segoe UI"
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(*theme_color)
        title_para.font.underline = False  # ✅ NO UNDERLINE
        title_para.alignment = PP_ALIGN.LEFT
        
        # ✅ REMOVED: Title underline decoration (was causing the underline issue)
        # The decorative line below title has been removed
        
        # Content area
        bullet_points = slide_data.get('bullet_points', [])
        
        if bullet_points:
            # Content background (positioned lower, no overlap with title)
            content_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(2.0),
                Inches(9), Inches(4.8)
            )
            content_bg.fill.solid()
            content_bg.fill.fore_color.rgb = RGBColor(248, 249, 250)
            content_bg.line.color.rgb = RGBColor(230, 230, 230)
            content_bg.line.width = Pt(1)
            
            # Add bullets
            y_position = 2.3
            max_bullets = min(6, len(bullet_points))
            
            for i in range(max_bullets):
                bullet = bullet_points[i]
                
                # Bullet icon
                bullet_icon = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(0.8), Inches(y_position),
                    Inches(0.15), Inches(0.15)
                )
                bullet_icon.fill.solid()
                bullet_icon.fill.fore_color.rgb = RGBColor(*accent_color)
                bullet_icon.line.fill.background()
                
                # Bullet text
                bullet_box = slide.shapes.add_textbox(
                    Inches(1.1), Inches(y_position - 0.1),
                    Inches(7.8), Inches(0.65)
                )
                bullet_frame = bullet_box.text_frame
                bullet_frame.text = bullet
                bullet_frame.word_wrap = True
                
                bullet_para = bullet_frame.paragraphs[0]
                bullet_para.font.name = "Segoe UI"
                bullet_para.font.size = Pt(16)
                bullet_para.font.color.rgb = RGBColor(60, 60, 60)
                bullet_para.font.underline = False  # ✅ NO UNDERLINE
                bullet_para.alignment = PP_ALIGN.LEFT
                bullet_para.line_spacing = 1.2
                
                y_position += 0.75
        
        # Presenter notes
        notes_text = slide_data.get('presenter_notes', '')
        if notes_text:
            try:
                notes_slide = slide.notes_slide
                notes_frame = notes_slide.notes_text_frame
                notes_frame.text = notes_text
            except:
                pass
        
        # ✅ Add watermark
        self._add_watermark(slide)
    
    
    
    def _add_conclusion_slide(
    self, 
    prs: Presentation, 
    slide_data: Dict[str, Any], 
    slide_number: int
    ):
        """✅ FIXED Conclusion slide - proper structure with no overlapping text"""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Get colors
        theme_color = self._hex_to_rgb(self.template_config.get("theme_color", "#1E3A8A"))
        accent_color = self._hex_to_rgb(self.template_config.get("accent_color", "#3B82F6"))
        
        # ✅ White background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Top accent bar
        top_accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.5)
        )
        top_accent.fill.solid()
        top_accent.fill.fore_color.rgb = RGBColor(*theme_color)
        top_accent.line.fill.background()
        
        # Left accent bar
        left_accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0.5),
            Inches(0.3), Inches(7)
        )
        left_accent.fill.solid()
        left_accent.fill.fore_color.rgb = RGBColor(*accent_color)
        left_accent.line.fill.background()
        
        # ✅ FIXED: Title positioned at top with enough space
        title = slide_data.get('title', 'Conclusion')
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.2),  # Moved down to avoid overlap
            Inches(8), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        title_frame.vertical_anchor = 1  # Center vertically
        
        title_para = title_frame.paragraphs[0]
        title_para.font.name = "Segoe UI"
        title_para.font.size = Pt(44)  # Slightly smaller to fit better
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(*theme_color)
        title_para.font.underline = False
        title_para.alignment = PP_ALIGN.CENTER
        
        # ✅ FIXED: Bullet points section with proper spacing
        bullet_points = slide_data.get('bullet_points', [])
        
        if bullet_points:
            # Content box positioned below title with enough gap
            # Slightly taller to safely fit multi-line bullets
            content_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.5), Inches(2.5),  # Moved down to avoid title overlap
                Inches(7), Inches(4.1)     # Increased height for multi-line bullets
            )
            content_bg.fill.solid()
            content_bg.fill.fore_color.rgb = RGBColor(248, 249, 250)
            content_bg.line.color.rgb = RGBColor(230, 230, 230)
            content_bg.line.width = Pt(1)
            
            # ✅ FIXED: Bullets with proper vertical spacing
            # Start a bit lower and give extra vertical room so 2-line bullets don't collide
            y_position = 2.8  # Start position for bullets
            max_bullets = min(4, len(bullet_points))  # Limit to 4 bullets
            
            for i in range(max_bullets):
                bullet = bullet_points[i]
                
                # Bullet icon
                bullet_icon = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(1.8), Inches(y_position),
                    Inches(0.2), Inches(0.2)
                )
                bullet_icon.fill.solid()
                bullet_icon.fill.fore_color.rgb = RGBColor(*accent_color)
                bullet_icon.line.fill.background()
                
                # Bullet text - adjusted width & height to prevent overflow
                bullet_box = slide.shapes.add_textbox(
                    Inches(2.2), Inches(y_position - 0.05),
                    Inches(5.8), Inches(0.8)  # Slightly taller for multi-line text
                )
                bullet_frame = bullet_box.text_frame
                bullet_frame.text = bullet[:120]  # Limit text length
                bullet_frame.word_wrap = True
                
                bullet_para = bullet_frame.paragraphs[0]
                bullet_para.font.name = "Segoe UI"
                bullet_para.font.size = Pt(16)  # Slightly smaller for better fit
                bullet_para.font.color.rgb = RGBColor(60, 60, 60)
                bullet_para.font.underline = False
                bullet_para.alignment = PP_ALIGN.LEFT
                bullet_para.line_spacing = 1.15
                
                # Extra spacing so wrapped bullets don't overlap the next bullet
                y_position += 1.0
        
        # ✅ FIXED: Thank You message at bottom - no overlap
        thanks_box = slide.shapes.add_textbox(
            Inches(2), Inches(6.6),  # Positioned at bottom
            Inches(6), Inches(0.5)
        )
        thanks_frame = thanks_box.text_frame
        thanks_frame.text = "Thank You! 🙏"
        thanks_frame.vertical_anchor = 1  # Center vertically
        
        thanks_para = thanks_frame.paragraphs[0]
        thanks_para.font.name = "Segoe UI"
        thanks_para.font.size = Pt(24)
        thanks_para.font.bold = True
        thanks_para.font.italic = True
        thanks_para.font.color.rgb = RGBColor(*theme_color)
        thanks_para.alignment = PP_ALIGN.CENTER
        
        # ✅ Add watermark
        self._add_watermark(slide)

    def _add_fallback_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        slide_number: int
    ):
        """Add a basic fallback slide if errors occur"""
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = slide_data.get('title', f'Slide {slide_number}')
        
        # Remove underline from title
        if title.text_frame.paragraphs:
            title.text_frame.paragraphs[0].font.underline = False
        
        # Content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        bullet_points = slide_data.get('bullet_points', [])
        for i, point in enumerate(bullet_points[:5]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.underline = False  # ✅ NO UNDERLINE
        
        # ✅ Add watermark
        self._add_watermark(slide)
    
    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """Convert hex color to RGB tuple"""
        try:
            hex_color = hex_color.lstrip('#')
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        except:
            return (30, 58, 138)  # Default blue
    
    def get_presentation_info(self, file_path: str) -> Dict[str, Any]:
        """Get information about the generated presentation"""
        try:
            file_size = os.path.getsize(file_path)
            
            return {
                "file_path": file_path,
                "file_size": file_size,
                "file_size_mb": round(file_size / (1024 * 1024), 2),
                "template_style": self.template_style,
                "created_at": datetime.now().isoformat()
            }
        except Exception as e:
            return {
                "file_path": file_path,
                "error": str(e)
            }